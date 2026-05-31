"""PPTX → markdown (markitdown) + structured slides (python-pptx).

markdown keeps markitdown's per-slide rendering. The structured form is
`{"slides": [{"slide": n, "text": str, "notes": str}]}` read directly with
python-pptx (already a markitdown dependency).
"""

from __future__ import annotations

from pathlib import Path

from markitdown.converters import PptxConverter

from ..config import settings
from ..image import ocr_embedded_images
from ..metadata import ooxml_core_props, pptx_slide_count
from .base import BaseParser, ParseResult
from .markitdown_util import convert_with_markitdown


class PptxParser(BaseParser):
    name = "markitdown-pptx"
    extensions = (".pptx",)
    isolation = True

    def parse(self, path: Path, mode: str) -> ParseResult:
        md = convert_with_markitdown(PptxConverter, ".pptx", path)
        section, ocr_stats = ocr_embedded_images(self._slide_images(path))
        if section:
            md = f"{md}\n\n{section}"
        structured = self._structured(path) if self.wants_structured(mode) else None
        metadata = ooxml_core_props(path)
        slides = pptx_slide_count(path)
        if slides is not None:
            metadata["slide_count"] = slides
        return ParseResult(
            parser=self.name,
            markdown=md,
            structured=structured,
            metadata=metadata,
            stats=ocr_stats,
        )

    @staticmethod
    def _slide_images(path: Path) -> list[tuple[str, bytes]]:
        """(label, bytes) for each picture shape, labeled by slide number.
        Empty unless embedded-image OCR is on (skips the python-pptx reopen)."""
        if not settings.ocr_embedded_images:
            return []
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        out: list[tuple[str, bytes]] = []
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                try:
                    out.append((f"Slide {i}", shape.image.blob))
                except Exception:
                    continue
        return out

    def _structured(self, path: Path) -> dict:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
            slides.append({"slide": i, "text": "\n".join(texts), "notes": notes})
        return {"slides": slides}
