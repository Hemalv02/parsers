"""Embedded-image OCR: text inside PPTX slide pictures and DOCX word/media
is recovered into a "## Embedded image text" section when enabled.

The parsers run in the SIGKILL worker subprocess (pptx) / in-process (docx);
the feature flag lives in settings, which a worker subprocess wouldn't see
monkeypatched — so these call the parsers directly (like test_pdf_routing),
with `ocr_embedded_images` toggled on the shared settings singleton. They
shell out to the real `tesseract` binary, like the existing image test."""

from __future__ import annotations

import io

import pytest


def _text_png(text: str, size: tuple[int, int] = (320, 80)) -> bytes:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", size, "white")
    ImageDraw.Draw(img).text((10, 30), text, fill="black")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _pptx_with_image(tmp_path, png: bytes):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(1), Inches(4), Inches(1))
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def _docx_with_image(tmp_path, png: bytes):
    import docx
    from docx.shared import Inches

    d = docx.Document()
    d.add_paragraph("Body paragraph.")
    d.add_picture(io.BytesIO(png), width=Inches(4))
    path = tmp_path / "doc.docx"
    d.save(str(path))
    return path


@pytest.fixture
def _ocr_on(monkeypatch):
    from app import config

    monkeypatch.setattr(config.settings, "ocr_embedded_images", True)
    monkeypatch.setattr(config.settings, "image_ocr_engine", "tesseract")


# --- PPTX -------------------------------------------------------------------


def test_pptx_embedded_image_ocr(tmp_path, _ocr_on):
    from app.parsers.pptx import PptxParser

    path = _pptx_with_image(tmp_path, _text_png("EMBEDDED SLIDE"))
    result = PptxParser().parse(path, "markdown")
    assert "## Embedded image text" in result.markdown
    assert "### Slide 1" in result.markdown
    assert "EMBEDDED" in result.markdown.upper()
    assert result.stats["images_ocred"] >= 1


def test_pptx_embedded_ocr_off_by_default(tmp_path):
    from app.parsers.pptx import PptxParser

    path = _pptx_with_image(tmp_path, _text_png("EMBEDDED SLIDE"))
    result = PptxParser().parse(path, "markdown")
    assert "## Embedded image text" not in result.markdown
    assert result.stats == {}  # no OCR work attempted


# --- DOCX -------------------------------------------------------------------


def test_docx_embedded_image_ocr(tmp_path, _ocr_on):
    from app.parsers.docx import DocxParser

    path = _docx_with_image(tmp_path, _text_png("DOCX PICTURE"))
    result = DocxParser().parse(path, "markdown")
    assert "## Embedded image text" in result.markdown
    assert "### image1.png" in result.markdown
    # OCR can drop a leading glyph on bitmap-rendered text; assert on a word
    # tesseract reads reliably rather than the exact string.
    assert "PICTURE" in result.markdown.upper()
    assert result.stats["images_ocred"] >= 1


# --- shared helper: dedup, size filter, stats -------------------------------


def test_ocr_embedded_images_dedup_and_size_filter(_ocr_on):
    from app.image import ocr_embedded_images

    big = _text_png("HELLO TEXT")  # 320x80, readable
    tiny = _text_png("x", size=(40, 40))  # below the 64px floor
    section, stats = ocr_embedded_images([("a.png", big), ("b.png", big), ("c.png", tiny)])
    assert stats["images_found"] == 3
    assert stats["images_ocred"] == 1  # b is a dup of a; c is too small
    assert stats["images_skipped"] == 2
    assert "## Embedded image text" in section


def test_ocr_embedded_images_noop_when_disabled():
    from app.image import ocr_embedded_images

    section, stats = ocr_embedded_images([("a.png", _text_png("HI THERE"))])
    assert section == ""
    assert stats == {}
