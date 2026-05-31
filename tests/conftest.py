"""Fixture factories: generate a small valid file per supported format.

Office files are generated with the same libraries the parsers use
(python-docx / python-pptx / openpyxl) so the round-trip is realistic.
The image fixture renders known text with Pillow so the tesseract OCR
assertion has something deterministic to find.
"""

from __future__ import annotations

import io
import json
from email.message import EmailMessage

import pytest
from fastapi.testclient import TestClient

from app.main import app


@pytest.fixture(scope="session")
def client() -> TestClient:
    return TestClient(app)


@pytest.fixture
def docx_bytes() -> bytes:
    import docx

    d = docx.Document()
    d.core_properties.title = "Quarterly Report"
    d.core_properties.author = "Finance Team"
    d.add_heading("Quarterly Report", level=1)
    d.add_paragraph("Revenue grew this quarter.")
    table = d.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    prs.core_properties.title = "Deck Title"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Deck Title"
    slide.placeholders[1].text = "First bullet point about parsing"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Total"])
    ws.append(["North", 42])
    ws.append(["South", 58])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture
def eml_bytes() -> bytes:
    msg = EmailMessage()
    msg["Subject"] = "Project kickoff"
    msg["From"] = "alice@example.com"
    msg["To"] = "bob@example.com"
    msg["Message-ID"] = "<abc123@example.com>"
    msg.set_content("Let's start the project on Monday.")
    return msg.as_bytes()


@pytest.fixture
def png_with_text() -> bytes:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (320, 80), "white")
    draw = ImageDraw.Draw(img)
    # Default bitmap font; large enough for tesseract to read reliably.
    draw.text((10, 30), "HELLO PARSER", fill="black")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def pdf_bytes() -> bytes:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_title("Annual Summary")
    pdf.set_author("Finance Team")
    pdf.add_page()
    pdf.set_font("Helvetica", size=14)
    pdf.cell(0, 10, "Annual Summary")
    pdf.ln(12)
    pdf.set_font("Helvetica", size=11)
    # The ported PDF path crops the top ~140px of page 1 (chapter-banner
    # suppression) and refuses pages below 100 text-chars/page (scanned-PDF
    # guard). So we lay down many lines that flow well past the top crop,
    # and spill onto a second page (which gets no header crop) to exercise
    # the normal per-page extraction + `## Page N` citation path.
    sentence = (
        "The parser extracts text per page and preserves page boundaries so "
        "a retrieval system can cite the exact page a passage came from."
    )
    for _ in range(15):
        pdf.multi_cell(0, 8, sentence, new_x="LMARGIN", new_y="NEXT")
    pdf.add_page()  # page 2 gets no header crop — guarantees dense content
    pdf.set_font("Helvetica", size=11)
    for _ in range(15):
        pdf.multi_cell(0, 8, sentence, new_x="LMARGIN", new_y="NEXT")
    out = pdf.output()  # fpdf2 returns a bytearray
    return bytes(out)


@pytest.fixture
def csv_bytes() -> bytes:
    return b"name,score\nalice,90\nbob,75\n"


@pytest.fixture
def json_bytes() -> bytes:
    return json.dumps({"project": "parsers", "files": ["a", "b"]}).encode()


@pytest.fixture
def html_bytes() -> bytes:
    return b"<html><body><h1>Title</h1><p>Hello <b>world</b></p></body></html>"


@pytest.fixture
def md_bytes() -> bytes:
    return b"# Heading\n\nSome **markdown** text.\n"


@pytest.fixture
def txt_bytes() -> bytes:
    return "Plain text with a unicode char: é\n".encode()
